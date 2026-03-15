import re
import PyPDF2
from docx import Document
from openpyxl import Workbook
import openpyxl
from youtube_transcript_api import YouTubeTranscriptApi
import requests
from bs4 import BeautifulSoup
from pptx import Presentation

def extract_text_from_pdf(pdf_file) -> str:
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                # Remove page numbers but preserve structural newlines
                page_text = re.sub(r'\n\s*\d+\s*\n', '\n', page_text)
                # Normalize multiple spaces but keep newlines
                page_text = re.sub(r'[ \t]+', ' ', page_text)
                text += page_text + "\n\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_docx(docx_file) -> str:
    try:
        doc = Document(docx_file)
        text = ""
        
        # Iterate over all elements maintaining order using doc.element.body
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph
        
        for element in doc.element.body:
            if isinstance(element, CT_P):
                p = Paragraph(element, doc)
                if p.text.strip():
                    text += p.text + "\n"
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                text += "\n"
                for i, row in enumerate(table.rows):
                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    text += "| " + " | ".join(row_data) + " |\n"
                    if i == 0:
                        text += "|" + "|".join(["---"] * len(row.cells)) + "|\n"
                text += "\n"
                
        return text.strip()
    except Exception as e:
        return f"Error extracting Word doc: {str(e)}"

def extract_text_from_xlsx(xlsx_file) -> str:
    try:
        wb = openpyxl.load_workbook(xlsx_file, data_only=True)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"\n--- Sheet: {sheet} ---\n"
            is_first_row = True
            for row in ws.iter_rows(values_only=True):
                # Preserving empty cells to maintain column alignment
                row_text = [str(cell) if cell is not None else "" for cell in row]
                # Only add the row if it's not completely empty
                if any(cell != "" for cell in row_text):
                    text += "| " + " | ".join(row_text) + " |\n"
                    if is_first_row:
                        text += "|" + "|".join(["---"] * len(row)) + "|\n"
                        is_first_row = False
        return text
    except Exception as e:
        return f"Error extracting Excel: {str(e)}"

def extract_text_from_txt(txt_file) -> str:
    try:
        # Expected to receive file contents as bytes in a file-like object
        return txt_file.read().decode("utf-8")
    except Exception as e:
        return f"Error extracting Text file: {str(e)}"

def extract_text_from_pptx(file) -> str:
    try:
        prs = Presentation(file)
        text = []
        filename = getattr(file, "filename", "unknown.pptx")
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.strip():
                        slide_content.append(shape.text.strip())
                elif shape.has_table:
                    slide_content.append("\n") # formatting space for table
                    for i, row in enumerate(shape.table.rows):
                        row_text = [cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells]
                        slide_content.append("| " + " | ".join(row_text) + " |")
                        if i == 0:
                            slide_content.append("|" + "|".join(["---"] * len(row.cells)) + "|")
                    slide_content.append("\n")
            if slide_content:
                text.append(f"--- [Slide {i+1}] ---\n" + "\n".join(slide_content))
        if not text:
            return f"POWERPOINT FILE ({filename}): [No text found or scanned images]"
        return f"POWERPOINT FILE ({filename}):\n" + "\n\n".join(text)
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_youtube_transcript(url: str) -> str:
    try:
        video_id = None
        patterns = [
            r'(?:v=|\/)([0-9A-Za-z_-]{11}).*',
            r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})',
            r'(?:embed\/)([0-9A-Za-z_-]{11})'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                video_id = match.group(1)
                break
        
        if not video_id:
            return "Error: Could not parse YouTube Video ID."

        try:
            # The library version in this environment uses instance methods
            api = YouTubeTranscriptApi()
            if hasattr(api, 'list'):
                transcript_list = api.list(video_id)
            else:
                # Fallback for standard versions
                transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        except Exception as trans_err:
             return f"Error listing transcripts: {str(trans_err)}"
        
        transcript = None
        try:
            transcript = transcript_list.find_transcript(['en'])
        except:
            pass
            
        if not transcript:
            try:
                transcript = transcript_list.find_generated_transcript(['en'])
            except:
                pass
        
        if not transcript:
            for t in transcript_list:
                transcript = t
                break
        
        if not transcript:
             return f"Error: No transcript available for {url}"

        data = transcript.fetch()
        
        # Robustly extract text from either dict-based or object-based snippets
        def get_text(item):
            if isinstance(item, dict):
                return item.get("text", "")
            return getattr(item, "text", "")

        text = " ".join([get_text(item) for item in data])
        return f"YOUTUBE TRANSCRIPT ({url}) [Language: {transcript.language}]:\n{text}"
    except Exception as e:
        err_msg = str(e)
        if "TranscriptsDisabled" in err_msg:
            return f"Error: Transcripts are disabled for this video ({url})."
        if "NoTranscriptFound" in err_msg:
            return f"Error: No transcript found for this video ({url})."
        return f"Error extracting YouTube transcript: {err_msg}"

def extract_text_from_url(url: str) -> str:
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        for script in soup(["script", "style"]):
            script.decompose()
            
        text = soup.get_text()
        
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return f"WEBSITE CONTENT ({url}):\n{text}"
    except Exception as e:
        return f"Error scraping URL: {str(e)}"
