"""
AI-Powered Instructional Design Automation System - Production Edition
========================================================================
Creates professional eLearning content using Groq's Llama 3.1 8B Instant

Key Features:
- Exact template matching for all outputs
- Professional Design Documents in Excel format
- Two storyboard formats (detailed and tabular)
- Powered by Groq Llama 3.1 8B Instant
- Robust error handling and validation
- ✨ NEW: Live inline editing for Design Doc & Storyboard
- ✨ NEW: AI Chat assistant for surgical targeted edits
"""

import streamlit as st
import os
from dotenv import load_dotenv
from groq import Groq
import PyPDF2
import io
import re
from datetime import datetime
from typing import Dict, Optional, List
import uuid
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
import json
import html
from youtube_transcript_api import YouTubeTranscriptApi
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
import glob

# ════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="AI Instructional Design System",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: 700;
        color: #1f4e78;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.5rem;
        font-weight: 600;
        color: #2c5282;
        margin-top: 2rem;
        margin-bottom: 1rem;
    }
    .info-box {
        background-color: #e6f3ff;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #1f4e78;
        margin: 1rem 0;
    }
    .success-box {
        background-color: #d4edda;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #28a745;
        margin: 1rem 0;
    }
    .edit-box {
        background-color: #fff8e1;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #FFC000;
        margin: 1rem 0;
    }
    .chat-box {
        background-color: #f0f4ff;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #4a6fa5;
        margin: 1rem 0;
    }
    .chat-user-msg {
        background-color: #e3f2fd;
        color: #000000;
        border-radius: 12px 12px 2px 12px;
        padding: 0.6rem 1rem;
        margin: 0.4rem 0;
        font-size: 0.93rem;
    }
    .chat-ai-msg {
        background-color: #f3e5f5;
        color: #000000;
        border-radius: 12px 12px 12px 2px;
        padding: 0.6rem 1rem;
        margin: 0.4rem 0;
        font-size: 0.93rem;
    }
    .stButton>button {
        width: 100%;
    }
    .tab-section {
        border: 1px solid #e0e0e0;
        border-radius: 8px;
        padding: 1rem;
        margin-top: 0.5rem;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
defaults = {
    'design_doc': None,
    'storyboard': None,
    'design_approved': False,
    'extracted_content': "",
    'intake_data': None,
    'storyboard_type': "Type 1",
    'intake_success': False,
    # ✨ NEW: edit states
    'design_doc_edited': None,       # holds live-edited version
    'storyboard_edited': None,       # holds live-edited storyboard
    'design_chat_history': [],       # chat messages for design doc AI
    'storyboard_chat_history': [],   # chat messages for storyboard AI
    'project_id': None               # Unique ID for auto-saving
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v





# ════════════════════════════════════════════════════════════════════════════
# UTILITY FUNCTIONS
# ════════════════════════════════════════════════════════════════════════════

HISTORY_DIR = "history"
if not os.path.exists(HISTORY_DIR):
    os.makedirs(HISTORY_DIR)

def init_project_id():
    """Ensure a valid project UUID exists."""
    if not st.session_state.get('project_id'):
        st.session_state['project_id'] = str(uuid.uuid4())

def save_project() -> Optional[str]:
    """Save current session state to a uniquely named JSON file."""
    try:
        init_project_id()
        project_id = st.session_state['project_id']
        filepath = os.path.join(HISTORY_DIR, f"{project_id}.json")
        
        # Determine title for meta-listing
        intake = st.session_state.get('intake_data')
        title = "Untitled"
        if intake and intake.get('course_title'):
             title = intake.get('course_title')
             
        # Add metadata useful for display
        metadata = {
             'title': title,
             'last_updated': datetime.now().isoformat()
        }
        
        keys_to_save = [
            'design_doc', 'storyboard', 'design_approved', 'extracted_content', 
            'intake_data', 'storyboard_type', 'intake_success', 'design_doc_edited',
            'storyboard_edited', 'design_chat_history', 'storyboard_chat_history',
            'project_id'
        ]
        
        data = {'_metadata': metadata}
        for k in keys_to_save:
            if k in st.session_state:
                data[k] = st.session_state[k]
                
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return filepath
    except Exception as e:
        print(f"Error saving project: {str(e)}") # Silent fail on auto-save
        return None

def load_project_from_file(filepath: str) -> bool:
    """Load session state from a JSON history file."""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            for k, v in data.items():
                if k != '_metadata':
                     st.session_state[k] = v
        return True
    except Exception as e:
        st.error(f"Error loading project from {filepath}: {str(e)}")
        return False

def get_saved_projects() -> List[Dict]:
    """Get list of saved projects sorted by newest first."""
    files = glob.glob(os.path.join(HISTORY_DIR, "*.json"))
    projects = []
    for f in files:
        try:
             with open(f, 'r', encoding='utf-8') as js:
                 data = json.load(js)
                 meta = data.get('_metadata', {})
                 projects.append({
                     'filepath': f,
                     'title': meta.get('title', 'Untitled'),
                     'last_updated': meta.get('last_updated', '1970-01-01T00:00:00')
                 })
        except:
             pass
    
    projects.sort(key=lambda x: x['last_updated'], reverse=True)
    return projects

def extract_text_from_pdf(pdf_file) -> str:
    """Extract clean text from uploaded PDF file."""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            page_text = re.sub(r'\n\s*\d+\s*\n', '\n', page_text)
            page_text = re.sub(r'\s+', ' ', page_text)
            text += page_text + "\n\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error extracting PDF: {str(e)}")
        return ""

def extract_text_from_docx(docx_file) -> str:
    """Extract text from Word document."""
    try:
        doc = Document(docx_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        st.error(f"Error extracting Word doc: {str(e)}")
        return ""

def extract_text_from_xlsx(xlsx_file) -> str:
    """Extract text from Excel file."""
    try:
        wb = Workbook()
        wb = openpyxl.load_workbook(xlsx_file, data_only=True)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"\n--- Sheet: {sheet} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    text += " | ".join(row_text) + "\n"
        return text
    except Exception as e:
        st.error(f"Error extracting Excel: {str(e)}")
        return ""

def extract_text_from_txt(txt_file) -> str:
    """Extract text from plain text file."""
    try:
        return txt_file.getvalue().decode("utf-8")
    except Exception as e:
        st.error(f"Error extracting Text file: {str(e)}")
        return ""

def extract_text_from_pptx(file) -> str:
    """Extract text from PowerPoint (.pptx) file."""
    try:
        prs = Presentation(file)
        text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                     if shape.text.strip():
                        slide_text.append(shape.text.strip())
            
            # Combine slide text
            if slide_text:
                text.append(f"[Slide {i+1}]: " + " | ".join(slide_text))
        
        if not text:
            return f"POWERPOINT FILE ({file.name}): [No text found or scanned images]"
            
        return f"POWERPOINT FILE ({file.name}):\n" + "\n\n".join(text)
            
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_youtube_transcript(url: str) -> str:
    """Extract transcript from YouTube video (robust)."""
    try:
        video_id = None
        # Robust regex for YouTube ID extraction
        patterns = [
            r'(?:v=|\/)([0-9A-Za-z_-]{11}).*',
            r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})',
            r'(?:embed\/)([0-9A-Za-z_-]{11})'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                video_id = match.group(1)
                break
        
        if not video_id:
            return "Error: Could not parse YouTube Video ID."

        # Instantiate the API (needed for newer versions/some environments)
        api = YouTubeTranscriptApi()
        
        # List all available transcripts
        try:
            transcript_list = api.list(video_id)
        except Exception as trans_err:
             return f"Error listing transcripts: {str(trans_err)}"
        
        transcript = None
        
        # 1. Try manual English
        try:
            transcript = transcript_list.find_transcript(['en'])
        except:
            pass
            
        # 2. Try generated English
        if not transcript:
            try:
                transcript = transcript_list.find_generated_transcript(['en'])
            except:
                pass
        
        # 3. Fallback: Take the first available one
        if not transcript:
            for t in transcript_list:
                transcript = t
                break
        
        if not transcript:
             return f"Error: No transcript available for {url}"

        data = transcript.fetch()
        text = " ".join([item.text for item in data])
        return f"YOUTUBE TRANSCRIPT ({url}) [Language: {transcript.language}]:\n{text}"

    except Exception as e:
        # Check for specific error messages to give better feedback
        err_msg = str(e)
        if "TranscriptsDisabled" in err_msg:
            return f"Error: Transcripts are disabled for this video ({url})."
        if "NoTranscriptFound" in err_msg:
            return f"Error: No transcript found for this video ({url})."
        return f"Error extracting YouTube transcript: {err_msg}"

def extract_text_from_url(url: str) -> str:
    """Extract text from a generic webpage."""
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Kill all script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
            
        text = soup.get_text()
        
        # Break into lines and remove leading/trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return f"WEBSITE CONTENT ({url}):\n{text}"
    except Exception as e:
        return f"Error scraping URL: {str(e)}"

def format_intake_text(intake_data: Dict) -> str:
    """Format intake data into readable text."""
    return f"""
Course Title: {intake_data['course_title']}
Business Unit: {intake_data['business_unit']}
Course Type: {intake_data['course_type']}
Target Audience: {intake_data['target_audience']}
Experience Level: {intake_data['experience_level']}
Geographic Spread: {intake_data['geographic_spread']}

Learning Objectives:
1. {intake_data['objective_1']}
2. {intake_data['objective_2']}
3. {intake_data['objective_3']}

Interactivity Level: {intake_data['interactivity_level']}
Output Required: {intake_data['output_required']}
"""

def get_active_design_doc() -> str:
    """Return the edited version if exists, else original."""
    return st.session_state.design_doc_edited or st.session_state.design_doc or ""

def get_active_storyboard() -> str:
    """Return the edited version if exists, else original."""
    return st.session_state.storyboard_edited or st.session_state.storyboard or ""

# ════════════════════════════════════════════════════════════════════════════
# STRATEGY DEFINITIONS
# ════════════════════════════════════════════════════════════════════════════

COURSE_LEVEL_STRATEGIES = {
    "Level 1": {
        "visual": """**LEVEL 1: Awareness-level (Text and primarily static images)**
Visual approach: Clean, simple, information-focused.
Recommended Strategies:
1. Text with supporting icons (Short bullets with relevant icons)
2. Infographics (Visual summaries of concepts or data)
3. Process flow diagrams (Simple step-by-step visuals)
4. Photographic visuals (Real images showing context or environment)
5. Charts and graphs (Data, trends, or comparisons)
6. Tables""",
        "interactivity": """**Level 1 Interaction Types (Basic):**
1. Click and Reveal (Learner clicks icons/labels -> Information appears)
2. Tabs Interaction (Content divided into tabs -> Learner explores each section)
3. Accordion Interaction (Expand/collapse sections -> Used for steps or rules)
4. Hotspots (Clickable areas on image -> Reveals information)""",
        "assessment": """**Level 1 Assessment Types (Knowledge Check):**
1. Multiple Choice (Single Correct)
2. Multiple Response (Multiple Correct)
3. True / False
4. Fill in the Blank
5. Matching (Terms and definitions, Process and outcomes)"""
    },
    "Level 2": {
        "visual": """**LEVEL 2: Engaging Courses (More contextual, learner-centered visuals)**
Includes all Level 1 visuals, plus:
1. Character-based illustrations (Workplace situations with characters)
2. Before-after comparison visuals (Correct vs incorrect behavior)
3. Demonstration videos (Task or process walkthroughs)
4. Conceptual illustrations (Visual metaphors for abstract ideas)
5. Simple animations
6. Expert talk (Recorded video of SME)""",
        "interactivity": """**Level 2 Interaction Types (Contextual Learning):**
Includes Level 1 interactions, plus:
1. Real-World News Incident (Short real-life case -> Shows consequences)
2. Mini Case Study (Short workplace scenario -> Reflection/Question)
3. Process Walkthrough (Step-by-step guided interaction)
4. Decision Point (Learner chooses an action -> Immediate feedback)""",
        "assessment": """**Level 2 Assessment Types (Applied Understanding):**
Includes Level 1 assessments, plus:
1. Scenario-Based MCQ (Short situation -> Learner chooses best action)
2. Sequencing / Ordering (Arrange steps in correct order)
3. Drag and Drop (Categorization: Do vs Don't, Risk vs Safe)"""
    },
    "Level 3": {
        "visual": """**LEVEL 3: Applied / Scenario-Based Courses (Highly contextual, realistic)**
Includes Level 1 & 2 visuals, plus:
1. Highly realistic scenarios
2. Complex simulations
3. Branching visual paths""",
        "interactivity": """**Level 3 Interaction Types (Advanced):**
Includes Level 1 & 2 interactions, plus:
1. Branching Scenario (Multiple decision paths -> Different endings)
2. Simulation (Software or Role-based)
3.    - Show me: System simulation or conversation simulation
4.    - Try me: Learner performs actions""",
        "assessment": """**Level 3 Assessment Types (Performance-based):**
Includes Level 1 & 2 assessments, plus:
1. Branching Scenario (Multiple decisions, different outcomes)
2. Simulation-Based Assessment (Learner performs actual steps)"""
    }
}

def get_strategy_for_level(interactivity_level: str) -> Dict[str, str]:
    if "Level 1" in interactivity_level:
        return COURSE_LEVEL_STRATEGIES["Level 1"]
    elif "Level 2" in interactivity_level:
        return COURSE_LEVEL_STRATEGIES["Level 2"]
    elif "Level 3" in interactivity_level or "Level 4" in interactivity_level:
        return COURSE_LEVEL_STRATEGIES["Level 3"]
    else:
        return COURSE_LEVEL_STRATEGIES["Level 2"]

# ════════════════════════════════════════════════════════════════════════════
# AI GENERATION FUNCTIONS USING GROQ
# ════════════════════════════════════════════════════════════════════════════

def generate_design_document(api_key: str, intake_data: Dict, content: str) -> Optional[str]:
    """Generate Design Document using Groq Llama 3.1 8B Instant."""
    try:
        client = Groq(api_key=api_key)
        strategies = get_strategy_for_level(intake_data['interactivity_level'])
        
        prompt = f"""You are an expert Instructional Designer creating a comprehensive Design Document.
        
INTAKE INFORMATION:
{format_intake_text(intake_data)}

REQUIRED VISUAL STRATEGIES:
{strategies['visual']}

REQUIRED INTERACTIVITY TYPES:
{strategies['interactivity']}

REQUIRED ASSESSMENT TYPES:
{strategies['assessment']}

SOURCE CONTENT:
{content[:15000]}

TASK:
Create a DETAILED Design Document following this EXACT structure.
The user wants "EXtreme Detail" and "Human Creativity" - do not summarize.

1. PROJECT INFORMATION
   - Project Name: {intake_data['course_title']}
   - Business Unit: {intake_data['business_unit']}
   - Target Audience: {intake_data['target_audience']}
   - Experience Level: {intake_data['experience_level']}
   - Date: {datetime.now().strftime('%Y-%m-%d')}

2. COURSE OVERVIEW
   - Context/Background: [2-3 sentences on why this training is needed based on the source content]
   - Project Goal: [1-2 sentences on business outcomes]
   - Duration: [Estimated total course duration]

3. LEARNING OBJECTIVES
   - {intake_data['objective_1']}
   - {intake_data['objective_2']}
   - {intake_data['objective_3']}

4. MODULE BREAKDOWN (CRITICAL: Use this EXACT table format with pipe separators - NO Markdown formatting inside cells)

Here is a GOLD STANDARD EXAMPLE of the detail required (Cybersecurity theme):
| Module | Delivery Mode | Learning Objectives | Topics | Recommended Strategy | Activities/Assessment | Duration |
|--------|---------------|---------------------|--------|----------------------|-----------------------|----------|
| Module 1: Introduction to Cybersecurity | Self-paced eLearning | • Define basic cyber security concepts.<br>• Explain why cyber security is critical for organizations.<br>• Identify common types of malware and cyber attacks encountered in day-to-day work.<br>• Recognize early warning signs of cyber threats.<br>• Understand the impact of security breaches on business continuity. | • Definition of Cyber Security<br>• Importance of Cyber Security<br>• Cyber Security in the Digital World<br>• Malware & Ransomware:<br>&nbsp;&nbsp;- Definition of malware and its types (viruses, worms, spyware, trojans).<br>&nbsp;&nbsp;- What is ransomware and how does it work?<br>&nbsp;&nbsp;- Consequences of malware and ransomware attacks.<br>• Social Engineering Attack:<br>&nbsp;&nbsp;- What is social engineering?<br>&nbsp;&nbsp;- Examples: pretexting, baiting, tailgating, and impersonation. | Real-world cyberattack examples set the stage, immediately immersing learners in the stakes of security. Interactive simulations then challenge users to identify phishing attempts in a safe, controlled environment. Case studies reveal the business impact of data breaches, followed by role-playing activities where learners must make critical security decisions under time pressure. The learning path concludes with a hands-on drag-and-drop exercise matching threats to defense strategies. | • Drag and drop cybersecurity concepts<br>• Multiple-choice questions on cybersecurity basics. | 2 hour |

NOW, generate the table for THIS course ({intake_data['course_title']}) following that EXACT LEVEL OF DETAIL.

CRITICAL INSTRUCTIONS FOR GENERATION:
1.  **MODULE COUNT**: You MUST generate exactly {intake_data['num_modules']} content modules (Module 1 to Module {intake_data['num_modules']}).
2.  **FORBIDDEN PHRASES**: 
    *   **NEVER start Objectives with**: "This module will...", "Learners will be able to...", "By the end of this module...". **Start directly with the verb** (e.g., "Analyze...", "Create...", "Identify...").
    *   **NEVER start Strategies with**: "This module will...", "In this module...", "Learners will...". **Start with the action** (e.g., "A simulation explores...", "Case studies highlight...", "Interactive scenarios guided the learner...").
3.  **VARIETY**: Every module MUST sound different. Do not repeat sentence structures.
4.  **ALIGNMENT**: Ensure "Recommended Strategy" and "Activities/Assessment" align with the specific Learning Objectives and ID Principles provided.

| Module | Delivery Mode | Learning Objectives | Topics | Recommended Strategy | Activities/Assessment | Duration |
|--------|---------------|---------------------|--------|----------------------|-----------------------|----------|
| Module 1: [Title] | Self-paced eLearning | • [Strong Verb] [Objective 1 - Detailed]<br>• [Strong Verb] [Objective 2 - Detailed]<br>• [Strong Verb] [Objective 3 - Detailed]<br>• [Strong Verb] [Objective 4 - Detailed] | • [Main Topic 1]<br>&nbsp;&nbsp;- [Sub-point 1]<br>&nbsp;&nbsp;- [Sub-point 2]<br>• [Main Topic 2]<br>&nbsp;&nbsp;- [Sub-point 1]<br>&nbsp;&nbsp;- [Sub-point 2]<br>• [Main Topic 3] | [EXTREMELY DETAILED strategy. Write 4-6 full sentences. Start with an action or description, NOT "Learners will". Tell a story of the learning experience.] | • [Specific Activity aligned to objectives]<br>• [Quiz details] | [Time] |
| Module 2: [Title] | Self-paced eLearning | • [Strong Verb] [Objective 1]<br>• [Strong Verb] [Objective 2]<br>• [Strong Verb] [Objective 3]<br>• [Strong Verb] [Objective 4] | • [Main Topic 1]<br>• [Main Topic 2 with detail breakdown]<br>• [Main Topic 3] | [EXTREMELY DETAILED strategy. Different opening style than Module 1. "A branching scenario allows..."] | • [Activity]<br>• [Quiz] | [Time] |
| ... [GENERATE EXACTLY {intake_data['num_modules']} MODULES TOTAL] ... |
| Knowledge Check | Self-paced eLearning | • Assess understanding | MCQs; Scenario-based questions | Quiz Format | Multiple-choice quiz | 30 min |
| Summary & Conclusion | Self-paced eLearning | • Review key concepts | Summary & Key takeaways | Recap points | Certificate of Completion | 15 min |

5. INSTRUCTIONAL STRATEGY
   - Pedagogy: [Approach based on {intake_data['interactivity_level']}]
   - Interactivity: [Specific interactive elements from REQUIRED INTERACTIVITY TYPES]
   - Media: [Specific visual strategies from REQUIRED VISUAL STRATEGIES]

6. ASSESSMENT STRATEGY
   - Formative: [Knowledge checks details]
   - Summative: [Final assessment details]
   - Criteria: [Pass/fail criteria]

7. TECHNICAL SPECIFICATIONS
   - LMS: SCORM 1.2
   - Devices: Desktop/Laptop, Tablet

IMPORTANT INSTRUCTIONS:
- **TONE**: Write in a **natural, professional human voice**. Avoid AI buzzwords like "delve", "comprehensive tapestry", "ensure", "foster". Use active voice.
- **NO REPETITION**: Do not repeat phrasing across modules. Make each strategy unique and specific to the content.
- **NARRATIVE FLOW**: In "Recommended Strategy", tell a story of how the learner experiences the module.
- EXTRACT EXTENSIVE DETAILS from the source content.
- NO GENERIC PLACEHOLDERS.
- Do NOT use bold ** or italic * formatting inside the table cells. Keep text clean.
- Use <br> for line breaks within table cells.
- Ensure the table structure is perfect.

Generate the complete Design Document now:"""

        with st.spinner("🎨 Generating Design Document with Groq Llama 3.1 8B..."):
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are an expert Instructional Designer who creates detailed, professional design documents based on source materials."},
                    {"role": "user", "content": prompt}
                ],
                model="llama-3.1-8b-instant",
                temperature=0.7,
                max_tokens=8000,
            )
            return chat_completion.choices[0].message.content
            
    except Exception as e:
        st.error(f"Error generating Design Document: {str(e)}")
        return None


def ai_edit_document(api_key: str, current_doc: str, user_instruction: str, doc_type: str = "Design Document", chat_history: List[Dict] = None) -> Optional[Dict[str, str]]:
    """
    AI assistant that can chat AND edit documents.
    Returns JSON: {"assistant_reply": str, "updated_document": str}
    """
    try:
        client = Groq(api_key=api_key)
        
        # Construct history context
        history_context = ""
        if chat_history:
            history_context = "\nRECENT CONVERSATION HISTORY (Use this for context/memory):\n"
            # Get last 10 messages for better context
            recent_msgs = chat_history[-10:] if len(chat_history) > 10 else chat_history
            for msg in recent_msgs:
                role = "Designer" if msg['role'] == 'user' else "AI"
                history_context += f"{role}: {msg['content']}\n"
            history_context += "\n"

        system_prompt = f"""You are a helpful, intelligent AI Instructional Design Assistant.
You have two goals:
1. **Chat & Assist**: Answer questions, greet the user, and be a helpful partner.
2. **Edit the Document**: Modify the {doc_type} based on the user's instructions.

**CRITICAL BEHAVIORAL PROTOCOL**:
1. **ANALYZE CONTEXT**: You must read the `RECENT CONVERSATION HISTORY` to understand the user's intent.
   - If the user says "Yes", "Okay", "Do it", or "Go ahead", they are **CONFIRMING** your previous proposal.
   - You MUST identify what you proposed in the last turn and **IMMEDIATELY EXECUTE IT** in the `updated_document`.
2. **ACTION ORIENTED**: Do not ask for permission effectively twice. If the user agrees, **DO THE WORK**.
3. **REAL INTELLIGENCE**: Do not rely on specific keywords. Understand the *flow* of the conversation.
   - Example Context: 
     - User: "Change X to Y."
     - AI: "I can do that, but maybe Z is better?"
     - User: "No, just do X."
   - Action: You must perform X immediately.

4. **SOCIAL/CLOSING PROTOCOL**: 
   - If the user's input is purely social (e.g., "Thanks", "Good job", "Ok bye") or a closing remark WITHOUT a new request:
     - `updated_document`: Return the `current_doc` exactly as is.
     - `assistant_reply`: A warm, brief closing (e.g., "You're welcome! Let me know if you need anything else.").
     - DO NOT say "I have updated the document" if you didn't change anything.

**OUTPUT**: 
   - `assistant_reply`: A natural confirmation of the action taken (e.g., "Done. I've updated the activities as discussed.").
   - `updated_document`: The actual modified text.

**CRITICAL OUTPUT FORMAT**:
You must ALWAYS return a valid JSON object with exactly these two keys:
1. "assistant_reply": Your conversational response to the user.
2. "updated_document": The FULL, COMPLETE content of the {doc_type}.
   - If you made edits, this is the new version.
   - If no edits were needed (just chat), return the CURRENT document exactly as is.

**EDITING RULES**:
- Only change what is requested.
- Keep table structures, markdown, and formatting intact.
- Use the CHAT HISTORY to understand context (e.g., "undo that", "make it shorter").
"""
        
        user_prompt = f"""Here is the CURRENT {doc_type}:
---START OF DOCUMENT---
{current_doc}
---END OF DOCUMENT---

{history_context}
Designer's LATEST instruction:
"{user_instruction}"

Respond in JSON format only:
{{
  "assistant_reply": "...",
  "updated_document": "..."
}}"""
        
        with st.spinner(f"🤖 AI is thinking..."):
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                model="llama-3.1-8b-instant",
                temperature=0.3,
                max_tokens=8000,
                response_format={"type": "json_object"}
            )
            response_json = json.loads(chat_completion.choices[0].message.content)
            return response_json
            
    except Exception as e:
        st.error(f"Error in AI assistant: {str(e)}")
        return None


def generate_storyboard(api_key: str, design_doc: str, intake_data: Dict, 
                       content: str, storyboard_type: str) -> Optional[str]:
    """Generate Storyboard using Groq Llama 3.1 8B Instant."""
    try:
        client = Groq(api_key=api_key)
        strategies = get_strategy_for_level(intake_data['interactivity_level'])
        
        if storyboard_type == "Type 1":
            prompt = f"""You are an expert eLearning Storyboard Developer creating a detailed, production-ready, block-based storyboard.
This storyboard is for professional Instructional Designers who will use it to build custom eLearning courses.
Every screen MUST contain enough detail that a developer can build the slide without guessing or making assumptions.

DESIGN DOCUMENT:
{design_doc[:30000]}

INTAKE INFORMATION:
{format_intake_text(intake_data)}

REQUIRED VISUAL STRATEGIES:
{strategies['visual']}

REQUIRED INTERACTIVITY TYPES:
{strategies['interactivity']}

SOURCE CONTENT:
{content[:15000]}

TASK:
Create a professional **Type 1 (Vertical/Block)** eLearning storyboard for EXACTLY {intake_data['num_modules']} MODULES.
Adhere strictly to the requested format for EVERY screen.
**SCREEN COUNT**: Generate a detailed broken-down storyboard with **5-10 screens per module**.

CRITICAL INSTRUCTION:
First, look at the 'MODULE BREAKDOWN' table in the Design Document above.
You MUST generate a storyboard section for ALL {intake_data['num_modules']} MODULES listed there.
Do not skip any modules. Match the titles exactly.

═══════════════════════════════════════════════════════════════
TABLE INTEGRITY & TOPIC RULES (READ CAREFULLY)
═══════════════════════════════════════════════════════════════

1. **STRICT TABLE FORMAT** — This is crucial for display:
   - EVERY table row MUST correspond to EXACTLY one raw line of text in your output.
   - **NO LITERAL NEWLINES** inside table cells. Use `<br>` for ALL line breaks and vertical spacing.
   - **MANDATORY PIPE `|`**: You MUST include the pipe separator between columns. If you omit it, the table breaks.

2. **TOPIC DEDUPLICATION**:
   - `Topic: [Name]` appears ONLY ONCE at the start of a new topic section. Do NOT repeat it for every screen.
   - **Screen [M.S] Title**: MUST be a **descriptive, unique title** for that specific slide's content (e.g., "M1.2 Title: The Impact of Social Engineering").
   - **NEVER** just repeat the Topic header as the Screen Title.

3. **OST (On-Screen Text)** — Each screen MUST have **3-5 balanced, substantive bullet points**.
   - Each bullet must be **1-2 full sentences** of real instructional content.
   - Balance the depth: informative but not overcrowded.

4. **Audio Narration** — Must be a **COMPLETE spoken script of 5-8 full sentences**.
   - The audio EXPANDS on the OST — it does NOT just repeat it.

5. **NO GENERIC PLACEHOLDERS** — NEVER use [Details], [Content], [Script], etc.

═══════════════════════════════════════════════════════════════

**CRITICAL FORMATTING RULES**:
1. **NO NEWLINES inside table cells**. Use `<br>` for line breaks. This is absolutely essential to prevent table breakage.
2. **COLUMN SEPARATION (CRITICAL FIX)**: You MUST place exactly one pipe character ` | ` between the Audio Narration and the Visual Instructions to separate the two columns. DO NOT merge them!
   - **Column 1**: Contains On-Screen Text and ends with the Audio Script.
   - **Column 2**: MUST start in a new cell after the pipe ` | ` with `**Graphics:**`.
   - Correct formatting MUST look like: `...end of audio script.] | **Graphics:**<br><br>[Visuals...`
3. **NO DUPLICATE HEADERS**: The `Module [Number]: [Title]` header MUST appear **ONLY ONCE** at the start of the module section.
4. **STRICT LAYOUT**: **DO NOT write any text or labels BETWEEN tables.**
5. **PLAIN TEXT TILES**: **DO NOT use bolding `**` for Module or Screen titles.**
6. **CONTINUOUS NUMBERING**: Screen titles MUST be numbered `Screen M.S` (e.g., Screen 1.1, 1.2... 1.8).

=============================================================================
Module [Number]: [Module Name from Design Doc]
Topic: [Topic Number and Name from Design Doc - Show ONLY at start of topic]
=============================================================================

Screen [M.S] Title: [Specific Descriptive Title Reflecting the Unique Slide Content]

| Screen Content & Audio | Visual Instructions |
| :--- | :--- |
| **OST (On-Screen Text):**<br><br>• [Substantive point 1 — informative sentence]<br>• [Substantive point 2 — key concept/definition]<br>• [Substantive point 3 — example or process step]<br>• [Substantive point 4 — additional detail]<br><br>**Audio:**<br><br>[5-8 sentence complete narration script. Natural, conversational tone. Expands on the OST.] | **Graphics:**<br><br>[Specific visual description: layout, characters, icons, colors.]<br><br>**Navigation:**<br><br>Click Next to continue.<br><br>**Additional Note to PG:**<br>[Specific production notes for developers.] |

Generate the complete storyboard for ALL modules now:"""

        else:  # Type 2 - Tabular format
            prompt = f"""You are an expert eLearning Storyboard Developer creating a detailed, production-ready tabular storyboard.
This storyboard is for professional Instructional Designers who will use it to build custom eLearning courses.
Every row MUST contain enough detail that a developer can build the slide without guessing or making assumptions.

DESIGN DOCUMENT:
{design_doc[:30000]}

INTAKE INFORMATION:
{format_intake_text(intake_data)}

REQUIRED VISUAL STRATEGIES:
{strategies['visual']}

REQUIRED INTERACTIVITY TYPES:
{strategies['interactivity']}

SOURCE CONTENT:
{content[:15000]}

TASK:
Create a professional **Type 2 (Tabular)** storyboard for EXACTLY {intake_data['num_modules']} MODULES.
**ROW COUNT**: Generate **5-10 rows (screens) per module**.

═══════════════════════════════════════════════════════════════
TABLE INTEGRITY & CONTENT RULES (READ CAREFULLY)
════════════════════════════════════════════════════════════════

1. **STRICT TABLE FORMAT** — This is crucial:
   - EVERY table row MUST correspond to EXACTLY one raw line of text in your output.
   - **NO LITERAL NEWLINES** inside table cells. Use `<br>` for ALL line breaks.
   - **MANDATORY PIPE `|`**: You MUST include the pipe separator between ALL columns.

2. **Topics Column**:
   - Provide a **specific sub-topic name** for each row.
   - Do NOT just repeat the Module title. Use the actual subject of the screen.

3. **On-screen text column**:
   - Each row MUST have **3-5 balanced, substantive bullet points**.
   - Use `<br>` for spacing within the cell.
   - BALANCE: Deep instructional value without overcrowding the row.

4. **Audio Narration column**:
   - Must be a **COMPLETE spoken script of 5-8 full sentences**.
   - Expands on the OST; does not repeat it.

5. **Visual Instructions/Developer Notes column**:
   - Be specific about graphics, layout, and interaction mechanics.

6. **NO GENERIC PLACEHOLDERS**.

═══════════════════════════════════════════════════════════════

**CRITICAL FORMATTING RULES:**
1. **NO NEWLINES inside table cells**. Use `<br>` tags for all line breaks.
2. Use `•` for bullet points.
3. All content for a single row must be on ONE line of text.
4. **NO STRAY TEXT**: Do not write headers or text outside the table rows.
5. **INTEGRATED SECTIONS**: Knowledge Checks and Summaries must be rows within the table.
6. **MANDATORY PIPES `|` (CRITICAL)**: Make absolutely sure every single row has EXACTLY the correct number of pipes separating the 7 columns. Never combine Visual Instructions and On-screen text.

**HEADERS MUST BE EXACTLY:**
| Section | Topics | Visual Instructions/Developer Notes | On-screen text | Audio Narration | Status | Actions required |

=============================================================================
COURSE: {intake_data['course_title']}
=============================================================================

MODULE [Number]: [Module Name]

| Section | Topics | Visual Instructions/Developer Notes | On-screen text | Audio Narration | Status | Actions required |
| :--- | :--- | :--- | :--- | :--- | :--- | :--- |
| Intro | [Specific Descriptive Sub-Topic Name] | • Full-screen titled layout with themed background<br>• Professional styled illustration related to the module topic<br>• Animated entrance: header slides from top, bullets fade in | **[Module Title]**<br>• [Substantive sentence introducing the specific topic]<br>• [Substantive sentence previewing key learning outcomes]<br>• [Substantive sentence on practical application] | Welcome to this section where we'll explore [topic] in depth. [Context from source material]. By the end of this module, you'll be able to [outcome]. Let's jump into our first key concept. | Draft | Slide design, animation setup |

[Generate 5-10 rows per module covering all topics. Every row must be a SINGLE LINE OF TEXT in the raw output.]

Generate the complete tabular storyboard for ALL {intake_data['num_modules']} modules now:"""

        with st.spinner(f"📋 Generating {storyboard_type} Storyboard with Groq Llama 3.1 8B..."):
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are an expert eLearning Storyboard Developer who creates detailed, screen-by-screen storyboards. You always base content on the provided materials and never use generic placeholders."},
                    {"role": "user", "content": prompt}
                ],
                model="llama-3.1-8b-instant",
                temperature=0.7,
                max_tokens=8000,
            )
            return chat_completion.choices[0].message.content
    
    except Exception as e:
        st.error(f"Error generating Storyboard: {str(e)}")
        return None

# ════════════════════════════════════════════════════════════════════════════
# EXPORT FUNCTIONS
# ════════════════════════════════════════════════════════════════════════════

def export_design_doc_to_xlsx(design_doc_text: str, intake_data: Dict) -> io.BytesIO:
    """Export Design Document to Excel matching the template exactly."""
    try:
        output = io.BytesIO()
        wb = Workbook()
        ws = wb.active
        ws.title = "Design Document"
        
        header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        yellow_fill = PatternFill(start_color="FFC000", end_color="FFC000", fill_type="solid")
        grey_fill = PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid")
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )
        
        ws.merge_cells('A1:G1')
        cell = ws['A1']
        cell.value = f"Cybersecurity Training Program For {intake_data.get('business_unit', 'Organization')}"
        cell.fill = header_fill
        cell.font = Font(bold=True, color="FFFFFF", size=14)
        cell.alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[1].height = 25
        
        ws['A2'] = "Trainers"
        ws['A2'].font = Font(bold=True, size=10, color="0000FF")
        ws['A3'] = "Instructor 1\nInstructor 2"
        ws['A3'].alignment = Alignment(wrap_text=True, vertical='top')
        ws.row_dimensions[3].height = 30
        
        ws.merge_cells('A4:G4')
        cell = ws['A4']
        cell.value = f"Design Document for {intake_data.get('course_title', 'Your Course')} — Self-paced eLearning"
        cell.fill = yellow_fill
        cell.font = Font(bold=True, size=11)
        cell.alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[4].height = 25
        
        headers = ["Module", "Delivery Mode", "Learning Objectives", "Topics", 
                   "Recommended Strategy", "Activities/Assessment", "Duration"]
        for col_num, header in enumerate(headers, 1):
            cell = ws.cell(row=5, column=col_num)
            cell.value = header
            cell.fill = grey_fill
            cell.font = Font(bold=True, size=10)
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            cell.border = thin_border
        ws.row_dimensions[5].height = 30
        
        row_num = 6
        lines = design_doc_text.split('\n')
        in_table = False
        
        for line in lines:
            line = line.strip()
            if '|' in line and 'Module' in line and 'Delivery' in line:
                in_table = True
                continue
            if in_table and '---' in line:
                continue
            if in_table and line.startswith('|'):
                cells = [c.strip() for c in line.split('|')]
                cells = [c for c in cells if c]
                if len(cells) >= 5:
                    for col_idx, content in enumerate(cells[:7], 1):
                        clean_content = html.unescape(content)
                        clean_content = clean_content.replace('<br>', '\n').replace('•', '•').replace('**', '').replace('__', '')
                        clean_content = re.sub(r'<[^>]+>', '', clean_content)
                        cell = ws.cell(row=row_num, column=col_idx)
                        cell.value = clean_content
                        cell.alignment = Alignment(vertical='top', wrap_text=True)
                        cell.border = thin_border
                    ws.row_dimensions[row_num].height = 80
                    row_num += 1
            elif in_table and not line:
                break
        
        column_configs = [
            ('A', 35), ('B', 20), ('C', 50), ('D', 50),
            ('E', 80), ('F', 40), ('G', 15)
        ]
        for col_letter, width in column_configs:
            ws.column_dimensions[col_letter].width = width
            
        for row in ws.iter_rows(min_row=6, max_row=row_num-1):
            max_height = 60
            for cell in row:
                cell.alignment = Alignment(vertical='top', wrap_text=True)
                if cell.value:
                    col_letter = cell.column_letter
                    col_width = ws.column_dimensions[col_letter].width
                    effective_width = col_width * 0.9
                    if effective_width > 0:
                        val_str = str(cell.value)
                        lines_wrapped = len(val_str) / effective_width
                        lines_newlines = val_str.count('\n')
                        total_lines = lines_wrapped + lines_newlines + 1
                        cell_height = total_lines * 15
                        if cell_height > max_height:
                            max_height = cell_height
            ws.row_dimensions[row[0].row].height = max_height
        
        wb.save(output)
        output.seek(0)
        return output
        
    except Exception as e:
        st.error(f"Error exporting to Excel: {str(e)}")
        return None

def export_storyboard_to_docx(storyboard_text: str, intake_data: Dict) -> io.BytesIO:
    """Export Storyboard to Word format with Table support."""
    try:
        output = io.BytesIO()
        doc = Document()
        
        title = doc.add_heading(f"eLearning Storyboard", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_heading(intake_data['course_title'], level=1)
        doc.add_paragraph(f"Business Unit: {intake_data['business_unit']}")
        doc.add_paragraph(f"Target Audience: {intake_data['target_audience']}")
        doc.add_paragraph(f"Interactivity: {intake_data['interactivity_level']}")
        doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
        doc.add_paragraph("")
        
        lines = storyboard_text.split('\n')
        table_lines = []
        in_table = False
        headers = []
        seen_headers = set()

        for line in lines:
            line_stripped = line.strip()
            
            if line_stripped.startswith('|'):
                if '---' in line_stripped:
                    continue
                if not in_table:
                    in_table = True
                    headers = [h.strip() for h in line_stripped.split('|') if h.strip()]
                    table_lines = []
                else:
                    cells = [c.strip() for c in line_stripped.split('|')]
                    if line_stripped.startswith('|'): cells.pop(0)
                    if line_stripped.endswith('|'): cells.pop()
                    table_lines.append(cells)
                continue
            
            if in_table and not line_stripped.startswith('|'):
                if headers and table_lines:
                    table = doc.add_table(rows=1, cols=len(headers))
                    table.style = 'Table Grid'
                    hdr_cells = table.rows[0].cells
                    for i, header in enumerate(headers):
                        if i < len(hdr_cells):
                            hdr_cells[i].text = header
                            hdr_cells[i].paragraphs[0].runs[0].bold = True
                    for row_data in table_lines:
                        row = table.add_row().cells
                        for i, cell_text in enumerate(row_data):
                            if i < len(row):
                                clean_text = html.unescape(cell_text).replace('<br>', '\n').replace('**', '')
                                row[i].text = clean_text
                in_table = False
                headers = []
                table_lines = []

            if line_stripped.startswith('===') or line_stripped.startswith('---'):
                continue
                
            # Clean artifacts just in case
            clean_line = line_stripped.replace('**', '').strip()
            
            if clean_line.upper().startswith('MODULE'):
                # Deduplicate
                if clean_line.upper() in seen_headers:
                    continue
                seen_headers.add(clean_line.upper())
                
                # Module Title: Blue Bold (Heading 1)
                p = doc.add_heading('', level=1)
                run = p.add_run(clean_line)
                run.bold = True
                run.font.color.rgb = DocxRGBColor(31, 78, 120) # Blue
            elif clean_line.upper().startswith('SCREEN'):
                # Deduplicate (optional, but good for safety)
                if clean_line.upper() in seen_headers:
                    continue
                seen_headers.add(clean_line.upper())
                
                # Screen Title: Bold (Heading 2)
                p = doc.add_heading('', level=2)
                run = p.add_run(clean_line)
                run.bold = True
                run.font.color.rgb = DocxRGBColor(0, 0, 0) # Black
            elif ':' in clean_line:
                # Key-Value Pair Check (Topic, Target Audience, etc.)
                parts = clean_line.split(':', 1)
                key = parts[0].strip()
                val = parts[1].strip() if len(parts) > 1 else ""
                
                # Allow standard keys even if long or mixed case
                is_header_key = (key.isupper() and len(key) < 40) or \
                                any(k in key for k in ["Topic", "Date", "Target", "Business", "Interactivity"])
                
                if is_header_key:
                    p = doc.add_paragraph()
                    run = p.add_run(key + ':')
                    run.bold = True
                    run.font.color.rgb = DocxRGBColor(31, 78, 120)
                    if val:
                        p.add_run(' ' + val)
                else:
                    doc.add_paragraph(clean_line)
            elif clean_line:
                doc.add_paragraph(clean_line)
        
        if in_table and headers and table_lines:
            table = doc.add_table(rows=1, cols=len(headers))
            table.style = 'Table Grid'
            hdr_cells = table.rows[0].cells
            for i, header in enumerate(headers):
                if i < len(hdr_cells):
                    hdr_cells[i].text = header
            for row_data in table_lines:
                row = table.add_row().cells
                for i, cell_text in enumerate(row_data):
                    if i < len(row):
                        row[i].text = html.unescape(cell_text).replace('<br>', '\n').replace('**', '')

        doc.save(output)
        output.seek(0)
        return output
        
    except Exception as e:
        output = io.BytesIO()
        doc = Document()
        doc.add_heading("Storyboard Export (Fallback)", 0)
        doc.add_paragraph(storyboard_text)
        doc.save(output)
        output.seek(0)
        return output

# ════════════════════════════════════════════════════════════════════════════
# ✨ NEW: EDIT & CHAT PANEL COMPONENTS
# ════════════════════════════════════════════════════════════════════════════

def render_document_studio(api_key: str, doc_key: str, edited_key: str, 
                         chat_key: str, doc_type: str):
    """
    Unified Document Studio: View, Edit, and Refine in one component.
    """
    # ── State Management ─────────────────────────────────────────────────────
    mode_key = f"studio_mode_{doc_key}"
    if mode_key not in st.session_state:
        st.session_state[mode_key] = "View"
        
    # Initialize chat with greeting if empty
    if chat_key in st.session_state and not st.session_state[chat_key]:
        st.session_state[chat_key] = [{"role": "assistant", "content": "Hi! How can I help you today?"}]
    elif chat_key not in st.session_state:
        st.session_state[chat_key] = [{"role": "assistant", "content": "Hi! How can I help you today?"}]
        
    # Get active content
    original_content = st.session_state[doc_key]
    edited_content = st.session_state[edited_key]
    active_content = edited_content if edited_content else original_content
    
    if not active_content:
        st.warning("No content to display.")
        return

    # ── Toolbar ──────────────────────────────────────────────────────────────
    st.markdown(f"### 🛠️ {doc_type} Studio")
    
    col_modes, col_status = st.columns([3, 2])
    
    with col_modes:
        mode = st.radio(
            "Studio Mode",
            ["View", "Split Editor", "AI Assistant"],
            horizontal=True,
            key=f"radio_{mode_key}",
            label_visibility="collapsed"
        )
        st.session_state[mode_key] = mode

    with col_status:
        if edited_content:
            st.markdown(f"**Status:** 🟡 Edited (Draft) | [Revert to Original]({st.session_state[doc_key]})", unsafe_allow_html=True)
            if st.button("🗑️ Discard Edits & Revert", key=f"revert_btn_{doc_key}", type="secondary"):
                st.session_state[edited_key] = None
                st.rerun()
        else:
            st.markdown("**Status:** 🟢 Original AI Generation")

    st.markdown("---")

    # ── VIEW MODE ────────────────────────────────────────────────────────────
    if mode == "View":
        st.markdown(active_content)

    # ── SPLIT EDITOR MODE ────────────────────────────────────────────────────
    elif mode == "Split Editor":
        st.markdown(
            "<div style='background-color: #fff8e1; padding: 10px; border-left: 4px solid #ffc107; margin-bottom: 20px;'>"
            "<strong>✏️ Live Editor:</strong> Edit Markdown on the left, see the Preview on the right. "
            "Changes are applied immediately."
            "</div>", 
            unsafe_allow_html=True
        )
        
        col_edit, col_preview = st.columns(2)
        
        with col_edit:
            st.markdown("#### 📝 Source")
            new_text = st.text_area(
                "Edit Content",
                value=active_content,
                height=800,
                key=f"editor_{doc_key}",
                label_visibility="collapsed"
            )
            # Auto-save changes to the edited_key
            if new_text != active_content:
                st.session_state[edited_key] = new_text
                st.rerun()
        
        with col_preview:
            st.markdown("#### 👁️ Live Preview")
            st.markdown(active_content)

    # ── AI ASSISTANT MODE ────────────────────────────────────────────────────
    elif mode == "AI Assistant":
        st.markdown(
            "<div style='background-color: #e3f2fd; padding: 10px; border-left: 4px solid #2196f3; margin-bottom: 20px;'>"
            "<strong>🤖 AI Refinement:</strong> Chat with the AI to request targeted changes. "
            "The AI will update the document based on your instructions."
            "</div>", 
            unsafe_allow_html=True
        )
        
        col_preview_l, col_chat_r = st.columns([3, 2])
        
        with col_preview_l:
            st.markdown("#### 📄 Current Document")
            with st.container(height=600, border=True):
                st.markdown(active_content)
                
        with col_chat_r:
            st.markdown("#### 💬 Chat Assistant")
            
            # Chat History Container
            history_container = st.container(height=400, border=True)
            with history_container:
                for msg in st.session_state[chat_key]:
                    role_icon = "🧑‍💻" if msg['role'] == 'user' else "🤖"
                    bg_color = "#f0f2f6" if msg['role'] == 'user' else "#e8f5e9"
                    # Force black text for visibility
                    st.markdown(
                        f"<div style='background-color: {bg_color}; color: #000000; padding: 10px; border-radius: 10px; margin-bottom: 10px;'>"
                        f"<strong>{role_icon} {'You' if msg['role'] == 'user' else 'AI'}:</strong><br>{msg['content']}"
                        f"</div>", 
                        unsafe_allow_html=True
                    )
            
            # Chat Input
            with st.form(key=f"chat_form_{doc_key}", clear_on_submit=True):
                user_instruction = st.text_area(
                    "Instruction",
                    placeholder="e.g., 'Make the tone of Module 1 more formal' or 'Add a quiz question about Phishing'",
                    height=100
                )
                
                cols_btn = st.columns([1, 1])
                with cols_btn[0]:
                    send = st.form_submit_button("🚀 Refine", type="primary")
                with cols_btn[1]:
                    clear = st.form_submit_button("🗑️ Clear Chat")
            
            if send and user_instruction:
                # Add user msg
                st.session_state[chat_key].append({"role": "user", "content": user_instruction})
                
                # Call AI

                response = ai_edit_document(
                    api_key, 
                    active_content, 
                    user_instruction, 
                    doc_type,
                    chat_history=st.session_state[chat_key]
                )
                
                if response:
                    reply = response.get("assistant_reply", "")
                    updated_doc = response.get("updated_document", active_content)
                    
                    # Clean preamble if any (fallback cleanup)
                    clean_doc = updated_doc.strip()
                    for prefix in ["---START OF DOCUMENT---"]:
                        if clean_doc.startswith(prefix):
                            clean_doc = clean_doc[len(prefix):].strip()
                    if clean_doc.endswith("---END OF DOCUMENT---"):
                         clean_doc = clean_doc[:-len("---END OF DOCUMENT---")].strip()
                    
                    # Add AI response to chat
                    timestamp = datetime.now().strftime("%H:%M:%S")
                    st.session_state[chat_key].append({"role": "assistant", "content": f"{reply} <br><span style='color:grey; font-size:0.8em'>({timestamp})</span>"})
                    
                    if clean_doc != active_content:
                        st.session_state[edited_key] = clean_doc
                        save_project() # Auto-save doc edit
                        st.success(f"✅ Document updated!")
                    else:
                        # Just a chat response, no doc update
                        save_project() # Auto-save chat history
                        pass
                        
                    st.rerun()
            
            if clear:
                st.session_state[chat_key] = []
                save_project() # Auto-save cleared chat
                st.rerun()


# ════════════════════════════════════════════════════════════════════════════
# MAIN APPLICATION
# ════════════════════════════════════════════════════════════════════════════

def main():
    """Main application logic."""
    
    st.markdown('<p class="main-header">📚 AI Instructional Design System</p>', unsafe_allow_html=True)
    st.markdown("**Professional eLearning Content Creation** | Powered by Groq Llama 3.1 8B Instant")
    
    load_dotenv()
    
    with st.sidebar:
        st.header("🗄️ Project History")
        
        if st.button("➕ New Project", type="primary", use_container_width=True):
            for k in list(st.session_state.keys()):
                if k != 'project_id': # keep track of keys, wait, actually we want to wipe project_id to generate a new one
                    del st.session_state[k]
            st.rerun()
            
        st.markdown("---")
        
        saved_projects = get_saved_projects()
        if saved_projects:
            st.markdown("**Recent Projects**")
            for p in saved_projects:
                # Add a button for each project
                btn_label = f"📁 {p['title']}"
                if st.button(btn_label, key=f"load_{p['filepath']}", use_container_width=True):
                    if load_project_from_file(p['filepath']):
                        st.success(f"Loaded {p['title']}")
                        st.rerun()
        else:
            st.info("No saved projects found.")
            
        st.markdown("---")
        st.header("⚙️ Configuration")
        
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
            api_key = st.text_input("Enter Groq API Key", type="password")
        
        if not api_key:
            st.error("❌ Please provide Groq API Key")
            st.stop()
        
        st.markdown("---")
        st.markdown("### 📋 Template Settings")
        
        storyboard_type = st.radio(
            "Storyboard Format",
            ["Type 1", "Type 2"],
            help="Choose template format for storyboard"
        )
        st.session_state.storyboard_type = "Type 1" if "Type 1" in storyboard_type else "Type 2"
        
        st.markdown("---")
        st.markdown("### 📖 Workflow")
        st.markdown("""
        1. ✍️ Complete intake form
        2. 📤 Upload source PDFs
        3. 🎨 Generate Design Doc
        4. ✏️ Edit / Chat refine
        5. ✅ Approve & Generate Storyboard
        6. ✏️ Edit / Chat refine storyboard
        7. 💾 Download files
        """)
        
        st.markdown("---")
        
        # ✨ Quick status panel
        st.markdown("### 📊 Document Status")
        if st.session_state.design_doc:
            if st.session_state.design_doc_edited:
                st.markdown("📄 Design Doc: 🟡 Edited")
            else:
                st.markdown("📄 Design Doc: 🟢 Generated")
        else:
            st.markdown("📄 Design Doc: ⚪ Not generated")
            
        if st.session_state.storyboard:
            if st.session_state.storyboard_edited:
                st.markdown("🎬 Storyboard: 🟡 Edited")
            else:
                st.markdown("🎬 Storyboard: 🟢 Generated")
        else:
            st.markdown("🎬 Storyboard: ⚪ Not generated")

        if st.session_state.design_doc_edited or st.session_state.storyboard_edited:
            st.markdown("---")
            if st.button("🔄 Reset All Edits"):
                st.session_state.design_doc_edited = None
                st.session_state.storyboard_edited = None
                st.session_state.design_chat_history = []
                st.session_state.storyboard_chat_history = []
                st.rerun()

        st.info("**Using Groq Llama 3.1 8B Instant** - Fast and efficient")
    
    # ═══════════════════════════════════════════════════════════════════════
    # STEP 1: INTAKE FORM
    # ═══════════════════════════════════════════════════════════════════════
    
    st.markdown('<p class="sub-header">📝 Step 1: Instructional Design Intake Form</p>', unsafe_allow_html=True)
    
    if st.session_state.get('intake_success'):
        st.markdown('<div class="success-box">✅ Intake form saved successfully!</div>', unsafe_allow_html=True)
        st.session_state.intake_success = False
    
    with st.form("intake_form"):
        st.subheader("Section 1: Basic Course Details")
        
        course_title = st.text_input("1. Course Title *", placeholder="e.g., Cybersecurity Fundamentals")
        
        col1, col2 = st.columns(2)
        with col1:
            business_unit = st.selectbox("2. Business Unit / Function *",
                ["Sales", "Operations", "Claims", "HR", "IT", "Compliance", "Finance", "Other"])
        with col2:
            course_type = st.selectbox("3. Course Type *",
                ["Regulatory / Compliance", "Product Training", "Process Training",
                 "Soft Skills", "Systems Training", "Technical Training", "Other"])
        
        st.markdown("---")
        st.subheader("Section 2: Audience Details")
        
        target_audience = st.text_input("1. Primary Audience Role *", placeholder="e.g., IT Security Professionals")
        
        col1, col2 = st.columns(2)
        with col1:
            experience_level = st.selectbox("2. Experience Level *",
                ["New to role", "0-2 years", "2-5 years", "5+ years"])
        with col2:
            geographic_spread = st.text_input("3. Geographic Spread", placeholder="e.g., Global, North America")
        
        st.markdown("---")
        st.subheader("Section 3: Course Design Requirements")
        
        st.markdown("**1. Primary Learning Objectives**")
        objective_1 = st.text_input("Objective 1 *", placeholder="Learners will be able to...")
        objective_2 = st.text_input("Objective 2 (Optional)", placeholder="Learners will be able to...")
        objective_3 = st.text_input("Objective 3 (Optional)", placeholder="Learners will be able to...")
        
        interactivity_level = st.selectbox("2. Level of Interactivity *", [
            "Level 1 – Informational (content + graphics + knowledge checks)",
            "Level 2 – Medium Interactivity (Level 1 + animations + interactions)",
            "Level 3 – High Interactivity (scenarios, simulations)",
            "Level 4 – Gamification"
        ])
        
        output_required = st.selectbox("3. Output Required *",
            ["Design Document", "Storyboard", "Both Design Document and Storyboard"])
        
        num_modules = st.slider("4. Number of Learning Modules", min_value=3, max_value=12, value=5,
            help="Select the number of modules to generate (excluding Intro/Summary)")
        
        submit_intake = st.form_submit_button("💾 Save Intake Form", type="primary")
    
    if submit_intake:
        if not all([course_title, business_unit, target_audience, objective_1]):
            st.error("⚠️ Please fill in all required fields marked with *")
        else:
            st.session_state.intake_data = {
                "course_title": course_title,
                "business_unit": business_unit,
                "course_type": course_type,
                "target_audience": target_audience,
                "experience_level": experience_level,
                "geographic_spread": geographic_spread or "Not specified",
                "objective_1": objective_1,
                "objective_2": objective_2 or "",
                "objective_3": objective_3 or "",
                "interactivity_level": interactivity_level,
                "output_required": output_required,
                "num_modules": num_modules
            }
            st.session_state.intake_success = True
            save_project() # Auto-save
            st.rerun()
    
    # ═══════════════════════════════════════════════════════════════════════
    # STEP 2: UPLOAD SOURCE DOCUMENTS
    # ═══════════════════════════════════════════════════════════════════════
    
    st.markdown('<p class="sub-header">📄 Step 2: Upload Source Documents</p>', unsafe_allow_html=True)
    
    col_upload, col_url = st.columns(2)
    
    with col_upload:
        uploaded_files = st.file_uploader(
            "Upload reference materials (PDF, Word, Excel, TXT, PPTX)",
            type=['pdf', 'docx', 'xlsx', 'txt', 'pptx', 'ppt'],  # Allow 'ppt' but it might fail without conversion logic
            accept_multiple_files=True,
            help="Upload one or more files to use as source content."
        )
        
    with col_url:
        external_links = st.text_area(
            "🔗 External Links (YouTube or Websites)",
            placeholder="Paste links here (one per line)...\nhttps://www.youtube.com/watch?v=...\nhttps://en.wikipedia.org/wiki/...",
            height=100,
            help="Add YouTube videos for transcript extraction or website URLs for scraping."
        )
    
    if uploaded_files or external_links.strip():
        file_count = len(uploaded_files) if uploaded_files else 0
        link_count = len([l for l in external_links.split('\n') if l.strip()])
        st.success(f"✅ {file_count} file(s) and {link_count} link(s) ready to process")
        
        with st.expander("📋 Verified Sources"):
            if uploaded_files:
                st.markdown("**Files:**")
                for file in uploaded_files:
                    st.write(f"• {file.name} ({file.type})")
            if external_links.strip():
                st.markdown("**Links:**")
                for link in external_links.split('\n'):
                    if link.strip():
                        st.write(f"• {link.strip()}")
        
        if st.button("📖 Extract All Content", type="primary"):
            all_content = ""
            total_items = (len(uploaded_files) if uploaded_files else 0) + (link_count if external_links.strip() else 0)
            progress_bar = st.progress(0)
            current_idx = 0
            
            # Process Files
            if uploaded_files:
                for file in uploaded_files:
                    with st.spinner(f"Processing file: {file.name}"):
                        content = ""
                        if file.name.endswith('.pdf'):
                            content = extract_text_from_pdf(file)
                        elif file.name.endswith('.docx'):
                            content = extract_text_from_docx(file)
                        elif file.name.endswith('.xlsx'):
                            content = extract_text_from_xlsx(file)
                        elif file.name.endswith('.txt'):
                            content = extract_text_from_txt(file)
                        elif file.name.endswith('.pptx') or file.name.endswith('.ppt'):
                            if file.name.endswith('.ppt'):
                                 st.warning(f"⚠️ .ppt files (binary) are not fully supported. Please convert {file.name} to .pptx for best results.")
                            content = extract_text_from_pptx(file)
                        elif file.name.endswith('.xlsx'):
                            content = extract_text_from_xlsx(file)
                        elif file.name.endswith('.txt'):
                            content = extract_text_from_txt(file)
                        elif file.name.endswith('.pptx') or file.name.endswith('.ppt'):
                            if file.name.endswith('.ppt'):
                                 st.warning(f"⚠️ .ppt files (binary) are not fully supported. Please convert {file.name} to .pptx for best results.")
                            content = extract_text_from_pptx(file)
                            
                        if content:
                            all_content += f"\n\n{'='*80}\n"
                            all_content += f"SOURCE FILE: {file.name}\n"
                            all_content += f"{'='*80}\n\n{content}"
                        
                        current_idx += 1
                        progress_bar.progress(current_idx / total_items)

            # Process Links
            if external_links.strip():
                for link in external_links.split('\n'):
                    url = link.strip()
                    if url:
                        with st.spinner(f"Processing link: {url}"):
                            content = ""
                            if "youtube.com" in url or "youtu.be" in url:
                                content = extract_youtube_transcript(url)
                            else:
                                content = extract_text_from_url(url)
                            
                            all_content += f"\n\n{'='*80}\n"
                            all_content += f"SOURCE LINK: {url}\n"
                            all_content += f"{'='*80}\n\n{content}"
                            
                            current_idx += 1
                            progress_bar.progress(current_idx / total_items)
            
            st.session_state.extracted_content = all_content
            save_project() # Auto-save
            st.markdown(f'<div class="success-box">✅ Successfully extracted {len(all_content):,} characters of source context.</div>', unsafe_allow_html=True)
            
            with st.expander("📄 Preview Extracted Content (first 3000 chars)"):
                st.text_area("Content Preview", all_content[:3000] + "...", height=300, disabled=True)
    
    # ═══════════════════════════════════════════════════════════════════════
    # STEP 3: GENERATE DESIGN DOCUMENT
    # ═══════════════════════════════════════════════════════════════════════
    
    st.markdown('<p class="sub-header">🎨 Step 3: Generate Design Document</p>', unsafe_allow_html=True)
    
    if not st.session_state.intake_data:
        st.info("👆 Please complete Step 1: Intake Form")
    elif not st.session_state.extracted_content:
        st.info("👆 Please complete Step 2: Upload & Extract Content")
    else:
        st.markdown('<div class="info-box">✅ Ready to generate Design Document</div>', unsafe_allow_html=True)
        
        if st.button("🚀 Generate Design Document", type="primary"):
            design_doc = generate_design_document(
                api_key,
                st.session_state.intake_data,
                st.session_state.extracted_content
            )
            if design_doc:
                st.session_state.design_doc = design_doc
                st.session_state.design_doc_edited = None  # Reset edits on regenerate
                st.session_state.design_chat_history = []
                st.session_state.design_approved = False
                save_project() # Auto-save
                st.success("✅ Design Document generated!")
                st.rerun()
    
    # Display Design Document
    if st.session_state.design_doc:
        # ── ✨ DOCUMENT STUDIO ────────────────────────────────────────────────
        render_document_studio(
            api_key=api_key,
            doc_key='design_doc',
            edited_key='design_doc_edited',
            chat_key='design_chat_history',
            doc_type="Design Document"
        )
        
        # ── Download Buttons ──────────────────────────────────────────────
        st.markdown("#### 💾 Download & Actions")
        active_doc = get_active_design_doc()
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.download_button(
                label="📥 Download as Text",
                data=active_doc,
                file_name=f"design_doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
                help="Downloads the current version (with your edits if any)"
            )
        
        with col2:
            xlsx_data = export_design_doc_to_xlsx(active_doc, st.session_state.intake_data)
            if xlsx_data:
                st.download_button(
                    label="📥 Download Excel",
                    data=xlsx_data,
                    file_name=f"design_doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    help="Downloads the current version (with your edits if any)"
                )
        
        with col3:
            if st.button("🔄 Regenerate Whole Doc", key="regen_design"):
                st.session_state.design_doc = None
                st.session_state.design_doc_edited = None
                st.session_state.design_chat_history = []
                st.session_state.design_approved = False
                st.rerun()
        
        # ── Approval ──────────────────────────────────────────────────────
        st.markdown("---")
        st.markdown("### ✅ Review & Approval")
        st.markdown("Once you're happy with the Design Document (original or edited), approve it to proceed to the storyboard.")
        
        col1, col2 = st.columns([1, 3])
        with col1:
            if st.button("✓ Approve Design Document", type="primary"):
                st.session_state.design_approved = True
                save_project() # Auto-save
                st.success("✅ Design Document approved! Proceed to Step 4.")
                st.rerun()
        with col2:
            if st.session_state.design_approved:
                st.markdown("🟢 **Design Document is APPROVED**")
    
    # ═══════════════════════════════════════════════════════════════════════
    # STEP 4: GENERATE STORYBOARD
    # ═══════════════════════════════════════════════════════════════════════
    
    st.markdown('<p class="sub-header">📋 Step 4: Generate Storyboard</p>', unsafe_allow_html=True)
    
    if not st.session_state.design_doc:
        st.info("👆 Please complete Step 3: Generate Design Document")
    elif not st.session_state.design_approved:
        st.warning("⚠️ Please approve the Design Document before generating storyboard")
    else:
        st.markdown(f'<div class="info-box">✅ Ready to generate {st.session_state.storyboard_type} Storyboard</div>', unsafe_allow_html=True)
        
        if st.button("🚀 Generate Storyboard", type="primary"):
            # Use the approved (possibly edited) design doc as the base
            base_doc = get_active_design_doc()
            storyboard = generate_storyboard(
                api_key,
                base_doc,
                st.session_state.intake_data,
                st.session_state.extracted_content,
                st.session_state.storyboard_type
            )
            if storyboard:
                st.session_state.storyboard = storyboard
                st.session_state.storyboard_edited = None  # Reset edits on regenerate
                st.session_state.storyboard_chat_history = []
                save_project() # Auto-save
                st.success("✅ Storyboard generated!")
                st.rerun()
    
    # Display Storyboard
    if st.session_state.storyboard:
        # ── ✨ DOCUMENT STUDIO ────────────────────────────────────────────────
        render_document_studio(
            api_key=api_key,
            doc_key='storyboard',
            edited_key='storyboard_edited',
            chat_key='storyboard_chat_history',
            doc_type="Storyboard"
        )
        
        # ── Download Buttons ──────────────────────────────────────────────
        st.markdown("#### 💾 Download & Actions")
        active_storyboard = get_active_storyboard()
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.download_button(
                label="📥 Download as Text",
                data=active_storyboard,
                file_name=f"storyboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
                help="Downloads the current version (with your edits if any)"
            )
        
        with col2:
            docx_data = export_storyboard_to_docx(active_storyboard, st.session_state.intake_data)
            if docx_data:
                st.download_button(
                    label="📥 Download Word",
                    data=docx_data,
                    file_name=f"storyboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    help="Downloads the current version (with your edits if any)"
                )
        
        with col3:
            if st.button("🔄 Regenerate Whole Storyboard", key="regen_storyboard"):
                st.session_state.storyboard = None
                st.session_state.storyboard_edited = None
                st.session_state.storyboard_chat_history = []
                st.rerun()
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #666; padding: 20px;'>
        <p><strong>AI Instructional Design System - Production Edition</strong></p>
        <p>Powered by Groq Llama 3.1 8B Instant | Built with Streamlit</p>
        <p>✨ Now with Live Editing & AI Chat for targeted, surgical document refinement</p>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()